from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

# Define file paths
BASE_DIR = Path(__file__).resolve().parent.parent
PPT_PATH = BASE_DIR / "report.pptx"  
CHART_IMG = BASE_DIR / "chart.png"   

def export_to_ppt():
    # Create new PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]   # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add chart image to slide (position + size in inches)
    slide.shapes.add_picture(str(CHART_IMG), Inches(1), Inches(1), Inches(8), Inches(5))

    # Remove old file if exists before saving
    if PPT_PATH.exists():
        PPT_PATH.unlink()
    prs.save(PPT_PATH)

    print(f"PowerPoint exported: {PPT_PATH}")
    return PPT_PATH
